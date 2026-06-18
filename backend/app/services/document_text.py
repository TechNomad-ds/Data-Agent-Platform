"""Lightweight text extraction for study materials."""
from pathlib import Path


def extract_document_text(file_path: Path, ext: str) -> str:
    ext = ext.lower()
    if ext == "pdf":
        try:
            import fitz
            doc = fitz.open(str(file_path))
            try:
                return "\n".join(page.get_text() for page in doc)
            finally:
                doc.close()
        except Exception:
            return file_path.read_text(encoding="utf-8", errors="ignore")

    if ext == "docx":
        try:
            from docx import Document
            doc = Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return file_path.read_text(encoding="utf-8", errors="ignore")

    if ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides: list[str] = []
            for idx, slide in enumerate(prs.slides, start=1):
                parts: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        parts.append(text)
                notes = getattr(slide, "notes_slide", None)
                if notes:
                    for shape in notes.shapes:
                        text = getattr(shape, "text", "")
                        if text:
                            parts.append(text)
                if parts:
                    slides.append(f"--- 第 {idx} 页 ---\n" + "\n".join(parts))
            return "\n\n".join(slides)
        except Exception:
            return ""

    return ""


def document_page_count(file_path: Path, ext: str, text: str = "") -> int:
    ext = ext.lower()
    if ext == "pdf":
        try:
            import fitz
            doc = fitz.open(str(file_path))
            try:
                return len(doc)
            finally:
                doc.close()
        except Exception:
            return max(1, len(text.splitlines()) // 40 + 1) if text else 0

    if ext == "docx":
        return max(1, len(text.splitlines()) // 30 + 1) if text else 0

    if ext == "pptx":
        try:
            from pptx import Presentation
            return len(Presentation(str(file_path)).slides)
        except Exception:
            return 0

    return 0
