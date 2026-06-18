#!/usr/bin/env python3
"""PKU deployment acceptance smoke test.

This script generates small course materials locally and uploads them through
the public API. It avoids LLM calls by default, so it can run before model
keys/quotas are fully exercised.
"""
from __future__ import annotations

import argparse
import io
import json
import os
import sys
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import httpx


DEFAULT_PASSWORD = "PkuSmokeTest123456"
DEFAULT_SMOKE_EMAIL = "pku_smoke_acceptance@example.com"
DEFAULT_SMOKE_USERNAME = "pku_smoke_acceptance"


def _python_runtime_info() -> dict[str, Any]:
    version = ".".join(map(str, sys.version_info[:3]))
    info: dict[str, Any] = {"python_version": version}
    if sys.version_info[:2] not in ((3, 11), (3, 12)):
        info["python_warning"] = (
            "PKU deployment is validated for Python 3.11/3.12; "
            "recreate backend/venv if this is a deployment run."
        )
    return info


def _write_report(report_path: str, report: dict[str, Any]) -> None:
    out = Path(report_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")


@dataclass
class SmokeContext:
    client: httpx.Client
    headers: dict[str, str]
    report: dict[str, Any]


def _require_status(resp: httpx.Response, expected: int | tuple[int, ...], step: str) -> None:
    allowed = expected if isinstance(expected, tuple) else (expected,)
    if resp.status_code not in allowed:
        try:
            detail = resp.json()
        except Exception:
            detail = resp.text[:500]
        raise RuntimeError(f"{step} failed: HTTP {resp.status_code}, {detail}")


def _docx_bytes(title: str, paragraphs: list[str]) -> bytes:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required. Install backend requirements first.") from exc

    doc = Document()
    doc.add_heading(title, level=1)
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(title: str, bullets: list[str]) -> bytes:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required. Install backend requirements first.") from exc

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text = "\n".join(bullets)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_bytes(sheets: dict[str, list[dict[str, Any]]]) -> bytes:
    try:
        import pandas as pd
    except ImportError as exc:
        raise RuntimeError("pandas/openpyxl are required. Install backend requirements first.") from exc

    buf = io.BytesIO()
    with pd.ExcelWriter(buf) as writer:
        for sheet_name, rows in sheets.items():
            pd.DataFrame(rows).to_excel(writer, sheet_name=sheet_name, index=False)
    return buf.getvalue()


def _materials() -> dict[str, list[tuple[str, bytes, str]]]:
    return {
        "计算机体系结构": [
            (
                "architecture_notes.docx",
                _docx_bytes(
                    "计算机体系结构复习笔记",
                    [
                        "流水线通过重叠指令阶段提高吞吐，但会产生结构、数据和控制冒险。",
                        "平均访存时间 AMAT = HitTime + MissRate * MissPenalty。",
                        "缓存局部性包括时间局部性和空间局部性。",
                    ],
                ),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ),
            (
                "architecture_slides.pptx",
                _pptx_bytes(
                    "计算机体系结构",
                    ["CPI 与流水线停顿", "Cache miss 对 AMAT 的影响", "乱序执行与指令级并行"],
                ),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            (
                "architecture_personal_note.md",
                "# 个人笔记\n我擅长 Python，但对硬件执行过程不熟。\n".encode("utf-8"),
                "text/markdown",
            ),
        ],
        "线性代数复习": [
            (
                "linear_algebra_notes.docx",
                _docx_bytes(
                    "线性代数复习",
                    [
                        "矩阵可以理解为线性变换，特征向量是在变换后方向不变的向量。",
                        "正交投影公式可以写作 proj_u(v) = (v·u / u·u)u。",
                    ],
                ),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ),
            (
                "linear_algebra_slides.pptx",
                _pptx_bytes("线性代数", ["向量空间", "基与维数", "特征值与对角化"]),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            (
                "linear_algebra_exercises.csv",
                "topic,difficulty\nprojection,medium\neigenvalue,hard\n".encode("utf-8"),
                "text/csv",
            ),
            (
                "online_course_ops_multisheet.xlsx",
                _xlsx_bytes({
                    "课程目录": [
                        {"course_id": "C001", "course_name": "线性代数基础", "category": "数学", "price": 199},
                        {"course_id": "C002", "course_name": "矩阵专题", "category": "数学", "price": 299},
                    ],
                    "报名记录": [
                        {"student_id": "S001", "course_id": "C001", "channel": "校内推荐"},
                        {"student_id": "S002", "course_id": "C002", "channel": "搜索"},
                    ],
                    "学习日志": [
                        {"student_id": "S001", "course_id": "C001", "completed": True, "minutes": 80},
                        {"student_id": "S002", "course_id": "C002", "completed": False, "minutes": 35},
                    ],
                }),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ),
        ],
        "现代文学导论": [
            (
                "literature_notes.docx",
                _docx_bytes(
                    "现代文学导论阅读笔记",
                    [
                        "鲁迅常以冷峻叙事揭示社会结构性问题。",
                        "沈从文更强调乡土经验、自然伦理和抒情性。",
                    ],
                ),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ),
            (
                "literature_slides.pptx",
                _pptx_bytes("现代文学导论", ["鲁迅", "沈从文", "乡土叙事与现代性"]),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            (
                "literature_summary.txt",
                "复习重点：叙事视角、历史语境、文学理论关键词。\n".encode("utf-8"),
                "text/plain",
            ),
        ],
    }


def register_or_login(client: httpx.Client, email: str, username: str, password: str) -> dict[str, str]:
    login = client.post("/api/auth/login", json={"email": email, "password": password})
    if login.status_code == 200:
        token = login.json()["access_token"]
        return {"Authorization": f"Bearer {token}"}
    if login.status_code == 429:
        _require_status(login, 200, "login")

    reg = client.post(
        "/api/auth/register",
        json={"email": email, "username": username, "password": password},
    )
    if reg.status_code not in (201, 400):
        _require_status(reg, 201, "register")

    login = client.post("/api/auth/login", json={"email": email, "password": password})
    _require_status(login, 200, "login")
    token = login.json()["access_token"]
    return {"Authorization": f"Bearer {token}"}


def space_name(base: str, run_label: str) -> str:
    return f"{base}-{run_label}" if run_label else base


def create_space(ctx: SmokeContext, name: str) -> str:
    resp = ctx.client.post("/api/data-spaces", headers=ctx.headers, json={"name": name, "description": "PKU smoke test"})
    if resp.status_code == 400 and "同名" in resp.text:
        listed = ctx.client.get("/api/data-spaces", headers=ctx.headers)
        _require_status(listed, 200, "list spaces after duplicate")
        for item in listed.json():
            if item["name"] == name:
                return item["id"]
    _require_status(resp, 201, f"create space {name}")
    return resp.json()["id"]


def upload_files(ctx: SmokeContext, space_id: str, items: list[tuple[str, bytes, str]], label: str) -> list[dict[str, Any]]:
    files = [("files", (name, io.BytesIO(content), mime)) for name, content, mime in items]
    resp = ctx.client.post(f"/api/data-spaces/{space_id}/upload", headers=ctx.headers, files=files, timeout=120)
    _require_status(resp, 201, f"upload {label}")
    uploaded = resp.json()
    if len(uploaded) != len(items):
        raise RuntimeError(f"upload {label} expected {len(items)} files, got {len(uploaded)}")
    return uploaded


def preview_file(ctx: SmokeContext, space_id: str, file_id: str, filename: str) -> None:
    resp = ctx.client.get(f"/api/data-spaces/{space_id}/files/{file_id}/preview", headers=ctx.headers)
    _require_status(resp, 200, f"preview {filename}")
    data = resp.json()
    if data.get("type") == "unsupported":
        raise RuntimeError(f"preview {filename} unsupported: {data}")
    if filename.endswith(".xlsx"):
        if data.get("type") != "workbook" or data.get("sheet_count", 0) < 3:
            raise RuntimeError(f"preview {filename} did not expose all Excel sheets: {data}")
        sheet_names = [s.get("name") for s in data.get("sheets", [])]
        for expected_sheet in ("课程目录", "报名记录", "学习日志"):
            if expected_sheet not in sheet_names:
                raise RuntimeError(f"preview {filename} missing sheet {expected_sheet}: {sheet_names}")
        return
    content = data.get("content", "")
    if filename.endswith(".pptx") and "计算机体系结构" not in content and "线性代数" not in content and "现代文学" not in content:
        raise RuntimeError(f"preview {filename} did not include expected slide text")


def wait_processing_ready(ctx: SmokeContext, space_id: str, label: str, timeout_seconds: int = 30) -> dict[str, Any]:
    deadline = time.monotonic() + timeout_seconds
    last_status: dict[str, Any] = {}
    while True:
        status = ctx.client.get(f"/api/data-spaces/{space_id}/processing-status", headers=ctx.headers)
        _require_status(status, 200, f"processing status {label}")
        last_status = status.json()
        if last_status.get("all_ready"):
            return last_status
        if last_status.get("error", 0) > 0:
            raise RuntimeError(f"processing status {label} has errors: {last_status}")
        if time.monotonic() >= deadline:
            raise RuntimeError(f"processing status {label} did not become ready: {last_status}")
        time.sleep(1)


def run(args: argparse.Namespace) -> dict[str, Any]:
    report: dict[str, Any] = {
        "base_url": args.base_url,
        "started_at": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "run_label": args.run_label,
        "checks": [],
        "spaces": {},
    }
    report.update(_python_runtime_info())
    client = httpx.Client(base_url=args.base_url.rstrip("/"), timeout=60)
    try:
        headers = register_or_login(client, args.email, args.username, args.password)
        ctx = SmokeContext(client=client, headers=headers, report=report)
        report["checks"].append("auth")

        materials = _materials()
        space_ids: dict[str, str] = {}
        uploads: dict[str, list[dict[str, Any]]] = {}

        # Interleaved course creation/upload mirrors the manual acceptance flow.
        for course in ("计算机体系结构", "现代文学导论", "线性代数复习"):
            sid = create_space(ctx, space_name(course, args.run_label))
            space_ids[course] = sid
            uploads[course] = upload_files(ctx, sid, materials[course][:2], f"{course} first batch")
            report["spaces"][course] = {"id": sid, "uploaded": [f["filename"] for f in uploads[course]]}

        # Rename data-space check.
        renamed = space_name("计算机体系结构-验收", args.run_label)
        rename_resp = client.put(
            f"/api/data-spaces/{space_ids['计算机体系结构']}",
            headers=headers,
            json={"name": renamed, "description": "已验证可修改名称"},
        )
        _require_status(rename_resp, 200, "rename data space")
        if rename_resp.json()["name"] != renamed:
            raise RuntimeError("rename response did not return updated name")
        report["checks"].append("rename_space")

        # Finish remaining files in a deliberately interleaved order.
        more_arch = upload_files(ctx, space_ids["计算机体系结构"], materials["计算机体系结构"][2:], "arch second batch")
        more_lit = upload_files(ctx, space_ids["现代文学导论"], materials["现代文学导论"][2:], "literature second batch")
        more_alg = upload_files(ctx, space_ids["线性代数复习"], materials["线性代数复习"][2:], "algebra second batch")
        uploads["计算机体系结构"].extend(more_arch)
        uploads["现代文学导论"].extend(more_lit)
        uploads["线性代数复习"].extend(more_alg)
        report["checks"].append("interleaved_uploads")

        # 7-file batch upload rate-limit regression check.
        batch_items = [
            (f"batch_note_{i}.md", f"# 批量上传测试 {i}\n这是一条验收笔记。\n".encode("utf-8"), "text/markdown")
            for i in range(1, 8)
        ]
        batch_uploaded = upload_files(ctx, space_ids["线性代数复习"], batch_items, "7 file batch")
        if len(batch_uploaded) != 7:
            raise RuntimeError("7 file batch upload did not return 7 records")
        report["checks"].append("seven_file_batch_upload")

        # Preview representative Word/PPTX/XLSX files.
        for course, files in uploads.items():
            for f in files:
                if f["filename"].endswith((".docx", ".pptx", ".xlsx")):
                    preview_file(ctx, space_ids[course], f["id"], f["filename"])
        report["checks"].append("document_preview")

        # Processing status endpoint should be usable after uploads.
        for course, sid in space_ids.items():
            report["spaces"][course]["processing_status"] = wait_processing_ready(ctx, sid, course)
        report["checks"].append("processing_status")

        # Conversation creation verifies space binding without consuming LLM credits.
        conv = client.post(
            "/api/chat/conversations",
            headers=headers,
            json={"data_space_id": space_ids["计算机体系结构"], "model_id": args.model_id, "title": "体系结构验收对话"},
        )
        _require_status(conv, 201, "create bound conversation")
        if conv.json()["data_space_id"] != space_ids["计算机体系结构"]:
            raise RuntimeError("bound conversation did not preserve data_space_id")

        general = client.post(
            "/api/chat/conversations",
            headers=headers,
            json={"data_space_id": None, "model_id": args.model_id, "title": "通用对话验收"},
        )
        _require_status(general, 201, "create general conversation")
        if general.json()["data_space_id"] is not None:
            raise RuntimeError("general conversation unexpectedly has data_space_id")
        report["checks"].append("conversation_space_binding")

        report["ok"] = True
        return report
    finally:
        client.close()


def main() -> int:
    parser = argparse.ArgumentParser(description="Run PKU deployment acceptance smoke checks.")
    parser.add_argument("--base-url", default=os.getenv("DATAMIND_BASE_URL", "http://127.0.0.1:8002"))
    parser.add_argument("--email", default=os.getenv("DATAMIND_SMOKE_EMAIL", DEFAULT_SMOKE_EMAIL))
    parser.add_argument("--username", default=os.getenv("DATAMIND_SMOKE_USERNAME", DEFAULT_SMOKE_USERNAME))
    parser.add_argument("--password", default=os.getenv("DATAMIND_SMOKE_PASSWORD", DEFAULT_PASSWORD))
    parser.add_argument("--model-id", default=os.getenv("DATAMIND_SMOKE_MODEL_ID", "smoke-test-model"))
    parser.add_argument("--run-label", default=os.getenv("DATAMIND_SMOKE_RUN_LABEL", time.strftime("%Y%m%d%H%M%S")))
    parser.add_argument("--report", default=os.getenv("DATAMIND_SMOKE_REPORT", "eval/pku_acceptance_report.json"))
    args = parser.parse_args()

    try:
        report = run(args)
    except Exception as exc:
        report = {
            "ok": False,
            "error": str(exc),
            "base_url": args.base_url,
            "ended_at": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        }
        report.update(_python_runtime_info())
        _write_report(args.report, report)
        print(json.dumps(report, ensure_ascii=False, indent=2), file=sys.stderr)
        return 1

    report["ended_at"] = time.strftime("%Y-%m-%dT%H:%M:%S%z")
    _write_report(args.report, report)
    print(json.dumps(report, ensure_ascii=False, indent=2))
    print(f"\nReport written to {args.report}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
